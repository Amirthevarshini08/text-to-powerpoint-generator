import streamlit as st
import openai
import json
import io
from pptx import Presentation

# --- Core Logic Functions ---

def get_slide_structure_from_llm(text, guidance, api_key):
    """
    Calls the LLM to get a JSON structure for the presentation slides.
    We'll use OpenAI for this example.
    """
    try:
        openai.api_key = api_key
        
        prompt = f"""
        Analyze the text below and structure it into presentation slides.
        The user's guidance is: "{guidance if guidance else 'A standard informative presentation'}".

        Return a valid JSON array of objects. Each object must have a "title" key and a "content" key.
        The "content" should be a list of strings (bullet points).
        Do not create a title slide or a thank you slide.

        TEXT:
        ---
        {text}
        ---
        """

        response = openai.chat.completions.create(
            model="gpt-3.5-turbo", # Using a faster, cheaper model for speed
            messages=[
                {"role": "system", "content": "You are an expert at structuring text into presentation slides. You only respond with the requested JSON."},
                {"role": "user", "content": prompt}
            ]
        )
        
        # The response content is a JSON string, so we need to parse it.
        raw_content = response.choices[0].message.content
        # A common issue: LLM sometimes wraps the JSON in ```json ... ```
        if raw_content.strip().startswith("```json"):
            raw_content = raw_content.strip()[7:-3]
            
        slide_data = json.loads(raw_content)
        return slide_data

    except Exception as e:
        # Provide a more specific error message to the user
        st.error(f"Error calling LLM: {e}. Check your API key and the input text.")
        return None

def create_presentation(structured_slides, template_file):
    """
    Generates a PowerPoint presentation from the structured slide data
    using the user's template.
    """
    try:
        # Load the user-provided template
        prs = Presentation(template_file)
        
        # --- Layout Selection Logic ---
        # Find the "Title and Content" layout. Its ID is usually 1, but we'll search for it.
        # This makes the app more robust to different templates.
        title_and_content_layout = None
        for i, layout in enumerate(prs.slide_layouts):
            # A simple heuristic to find the most common content slide
            if "Title" in layout.name and "Content" in layout.name:
                 title_and_content_layout = layout
                 break
        
        # Fallback if a specific layout isn't found
        if not title_and_content_layout:
            title_and_content_layout = prs.slide_layouts[1] # A common default

        # --- Slide Generation Loop ---
        for slide_data in structured_slides:
            slide = prs.slides.add_slide(title_and_content_layout)
            
            # Set the title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get("title", "")

            # Set the content (body)
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx >= 1: # Find the main body placeholder
                    body_shape = shape
                    break
            
            if body_shape:
                tf = body_shape.text_frame
                tf.clear() # Clear existing placeholder text
                
                content_points = slide_data.get("content", [])
                if isinstance(content_points, list) and content_points:
                    tf.text = content_points[0]
                    for point in content_points[1:]:
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 1

        # Save the presentation to a memory buffer to be downloaded
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

    except Exception as e:
        st.error(f"Error during presentation creation: {e}. Please check your template file.")
        return None


# --- Streamlit User Interface ---

st.set_page_config(layout="wide", page_title="Text to Presentation Generator")
st.title("✨ Your Text, Your Style – Auto-Generator")

# App description
st.markdown("Turn bulk text, markdown, or prose into a fully formatted PowerPoint presentation that matches your chosen template's look and feel.")

# Input fields arranged in columns
col1, col2 = st.columns(2)

with col1:
    st.subheader("1. Paste Your Text")
    input_text = st.text_area("Enter the full text for your presentation here...", height=350)
    
    st.subheader("2. Add Guidance (Optional)")
    guidance_text = st.text_input("e.g., 'Make it an investor pitch deck', 'A technical deep-dive'")

with col2:
    st.subheader("3. Provide API Key & Template")
    api_key = st.text_input("Your LLM API Key (OpenAI)", type="password", help="Your key is used only for this request and is not stored.")
    
    uploaded_template = st.file_uploader(
        "Upload a .pptx or .potx template",
        type=['pptx', 'potx']
    )
    
    st.markdown("---")
    generate_button = st.button("🚀 Generate Presentation", use_container_width=True)

# --- Main Logic on Button Click ---

if generate_button:
    # Validate inputs
    if not input_text:
        st.warning("Please paste some text to begin.")
    elif not api_key:
        st.warning("Please enter your OpenAI API key.")
    elif not uploaded_template:
        st.warning("Please upload a PowerPoint template.")
    else:
        with st.spinner("Analyzing text and building your presentation... This can take a minute."):
            # Step 1: Get structured data from LLM
            structured_slides = get_slide_structure_from_llm(input_text, guidance_text, api_key)
            
            if structured_slides:
                # Step 2: Create the presentation
                presentation_buffer = create_presentation(structured_slides, uploaded_template)
                
                if presentation_buffer:
                    st.success("🎉 Your presentation is ready!")
                    
                    # Step 3: Provide download link
                    st.download_button(
                        label="📥 Download Presentation",
                        data=presentation_buffer,
                        file_name="generated_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
